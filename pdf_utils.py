from typing import List
import re
from PyPDF2 import PdfReader
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import docx
except Exception:
    docx = None

def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    texts = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(texts)

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split `text` into semantically coherent, overlapping chunks.

    Uses a sentence-aware splitter and accumulates sentences until reaching
    `chunk_size` characters. Keeps an `overlap` (characters) between chunks to
    preserve context across boundaries. If a single sentence is larger than
    `chunk_size`, it will be split with a sliding window fallback.
    """
    if not text:
        return []

    # Normalize whitespace
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r"\s+", " ", text).strip()

    # Simple sentence splitter (works reasonably without heavy deps)
    sentences = re.split(r'(?<=[\.!?])\s+', text)

    chunks: List[str] = []
    cur = ""
    for s in sentences:
        s = s.strip()
        if not s:
            continue
        if len(cur) + len(s) + 1 <= chunk_size:
            cur = (cur + " " + s).strip() if cur else s
            continue

        # cur + s would exceed chunk_size
        if cur:
            chunks.append(cur)
            # prepare next chunk starting with an overlap of previous chunk's tail
            tail = chunks[-1][-overlap:] if overlap and len(chunks[-1]) >= overlap else ""
            cur = (tail + " " + s).strip()
        else:
            # current sentence alone is larger than chunk_size; break it down
            piece = s
            step = max(1, chunk_size - overlap)
            i = 0
            while i < len(piece):
                chunks.append(piece[i:i + chunk_size])
                i += step
            cur = ""

    if cur:
        chunks.append(cur)

    # If result is still empty for some reason, fallback to fixed windows
    if not chunks:
        s = text
        step = max(1, chunk_size - overlap)
        i = 0
        while i < len(s):
            chunks.append(s[i:i + chunk_size])
            i += step

    return chunks


def extract_text_from_pptx(path: str) -> str:
    if Presentation is None:
        raise RuntimeError('python-pptx not installed')
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                # Text frames (standard text boxes / placeholders)
                if getattr(shape, "has_text_frame", False):
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        runs = [r.text for r in para.runs if getattr(r, "text", None)]
                        if runs:
                            texts.append("".join(runs))
                    continue

                # Tables
                table = getattr(shape, "table", None)
                if table is not None:
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = getattr(cell, "text", None)
                            if cell_text:
                                texts.append(cell_text)
                    continue

                # Fallback to generic text attribute
                t = getattr(shape, "text", None)
                if t:
                    texts.append(t)
            except Exception:
                # ignore shapes we can't read
                continue
    return "\n".join(texts)


def extract_text_from_docx(path: str) -> str:
    if docx is None:
        raise RuntimeError('python-docx not installed')
    doc = docx.Document(path)
    texts = [p.text for p in doc.paragraphs if p.text]
    return "\n".join(texts)


def extract_text_from_file(path: str) -> str:
    """Dispatch to the right extractor based on file suffix.

    Supports: PDF, PPTX, DOCX. Raises for unsupported types.
    """
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == '.pdf':
        return extract_text_from_pdf(path)
    if suffix in ('.pptx',):
        return extract_text_from_pptx(path)
    if suffix in ('.docx',):
        return extract_text_from_docx(path)
    raise ValueError(f'Unsupported file type: {suffix}')
